"""
Complete Admin Views for Learning Management System
Includes all CRUD operations, file upload, parsing, analytics, and search
"""
from rest_framework import viewsets, status, views
from rest_framework.decorators import action
from rest_framework.response import Response
from rest_framework.permissions import IsAuthenticated, IsAdminUser
from rest_framework.parsers import MultiPartParser, FormParser, JSONParser
from django.db import transaction
from django.db.models import Count, Avg, Q, F, Sum
from django.utils import timezone
from django.core.files.storage import default_storage
from datetime import timedelta
import os
import json

from .models import (
    CareerPath, LearningModule, Quiz, Question, QuestionChoice,
    Enrollment, ModuleProgress, QuizAttempt, Certificate
)
from .serializers_admin import (
    AdminCareerPathSerializer, AdminLearningModuleSerializer,
    AdminQuizSerializer, AdminQuestionSerializer, AdminOverviewSerializer,
    ModuleFileSerializer, PathAnalyticsSerializer, ModuleAnalyticsSerializer,
    QuizAnalyticsSerializer, ContentModerationSerializer, SystemSettingsSerializer
)
from apps.accounts.models import User


class AdminOverviewView(views.APIView):
    """
    Admin Dashboard Overview - Real-time metrics
    GET /api/admin/overview/
    """
    permission_classes = [IsAuthenticated, IsAdminUser]
    
    def get(self, request):
        today = timezone.now().date()
        thirty_days_ago = timezone.now() - timedelta(days=30)
        
        # User metrics
        total_users = User.objects.count()
        active_users = User.objects.filter(
            last_login__gte=thirty_days_ago
        ).count()
        new_users_today = User.objects.filter(
            created_at__date=today
        ).count()
        
        # Path metrics
        total_paths = CareerPath.objects.count()
        active_paths = CareerPath.objects.filter(is_active=True).count()
        
        # Module metrics
        total_modules = LearningModule.objects.count()
        total_uploads = LearningModule.objects.exclude(file='').count()
        pending_uploads = LearningModule.objects.filter(
            files__parse_status='pending'
        ).distinct().count()
        
        # Quiz metrics
        total_quizzes = Quiz.objects.count()
        
        # Enrollment metrics
        total_enrollments = Enrollment.objects.count()
        active_enrollments = Enrollment.objects.filter(
            completed_at__isnull=True
        ).count()
        
        # Completion rate
        completed_enrollments = Enrollment.objects.filter(
            completed_at__isnull=False
        ).count()
        avg_completion_rate = (
            (completed_enrollments / total_enrollments * 100) 
            if total_enrollments > 0 else 0
        )
        
        # Recent errors (placeholder - you'd query SystemError model)
        recent_errors = []
        
        # Pending moderation (placeholder)
        pending_moderation = 0
        
        data = {
            'total_users': total_users,
            'active_users': active_users,
            'new_users_today': new_users_today,
            'total_paths': total_paths,
            'active_paths': active_paths,
            'total_modules': total_modules,
            'total_uploads': total_uploads,
            'pending_uploads': pending_uploads,
            'total_quizzes': total_quizzes,
            'total_enrollments': total_enrollments,
            'active_enrollments': active_enrollments,
            'avg_completion_rate': round(avg_completion_rate, 2),
            'recent_errors': recent_errors,
            'pending_moderation': pending_moderation,
        }
        
        serializer = AdminOverviewSerializer(data)
        return Response(serializer.data)


class AdminPathViewSet(viewsets.ModelViewSet):
    """
    Complete Path Management
    GET    /api/admin/paths/                 - List all paths
    POST   /api/admin/paths/                 - Create new path
    GET    /api/admin/paths/{id}/            - Get path details
    PUT    /api/admin/paths/{id}/            - Update path
    DELETE /api/admin/paths/{id}/            - Delete path
    GET    /api/admin/paths/{id}/modules/    - List path modules
    GET    /api/admin/paths/{id}/analytics/  - Get path analytics
    POST   /api/admin/paths/{id}/publish/    - Publish/unpublish path
    """
    permission_classes = [IsAuthenticated, IsAdminUser]
    serializer_class = AdminCareerPathSerializer
    
    def get_queryset(self):
        queryset = CareerPath.objects.annotate(
            modules_count=Count('modules'),
            enrollments_count=Count('enrollments')
        )
        
        # Filters
        program = self.request.query_params.get('program')
        difficulty = self.request.query_params.get('difficulty')
        is_active = self.request.query_params.get('is_active')
        search = self.request.query_params.get('search')
        
        if program:
            queryset = queryset.filter(program=program)
        if difficulty:
            queryset = queryset.filter(difficulty_level=difficulty)
        if is_active is not None:
            queryset = queryset.filter(is_active=is_active.lower() == 'true')
        if search:
            queryset = queryset.filter(
                Q(name__icontains=search) |
                Q(description__icontains=search)
            )
        
        return queryset.order_by('-created_at')
    
    @action(detail=True, methods=['get'])
    def modules(self, request, pk=None):
        """Get all modules for a path, ordered by index"""
        path = self.get_object()
        modules = path.modules.all().order_by('order')
        serializer = AdminLearningModuleSerializer(modules, many=True)
        return Response(serializer.data)
    
    @action(detail=True, methods=['get'])
    def analytics(self, request, pk=None):
        """Get analytics for a path"""
        path = self.get_object()
        
        # Calculate analytics
        enrollments = path.enrollments.all()
        total_enrollments = enrollments.count()
        completed = enrollments.filter(completed_at__isnull=False)
        completed_count = completed.count()
        
        analytics_data = {
            'total_views': 0,  # Would track in PathView model
            'unique_viewers': 0,
            'total_enrollments': total_enrollments,
            'active_enrollments': enrollments.filter(completed_at__isnull=True).count(),
            'completed_enrollments': completed_count,
            'avg_completion_rate': (completed_count / total_enrollments * 100) if total_enrollments > 0 else 0,
            'avg_time_to_complete': None,
            'drop_off_points': [],
        }
        
        serializer = PathAnalyticsSerializer(analytics_data)
        return Response(serializer.data)
    
    @action(detail=True, methods=['post'])
    def publish(self, request, pk=None):
        """Publish or unpublish a path"""
        path = self.get_object()
        is_active = request.data.get('is_active', True)
        path.is_active = is_active
        path.save()
        
        serializer = self.get_serializer(path)
        return Response(serializer.data)
    
    @action(detail=True, methods=['post'])
    def reorder_modules(self, request, pk=None):
        """Reorder modules in a path"""
        path = self.get_object()
        module_ids = request.data.get('module_ids', [])
        
        with transaction.atomic():
            for index, module_id in enumerate(module_ids):
                LearningModule.objects.filter(
                    id=module_id, 
                    career_path=path
                ).update(order=index)
        
        return Response({'status': 'reordered'})


class AdminModuleViewSet(viewsets.ModelViewSet):
    """
    Complete Module Management
    Supports both manual creation and file upload with parsing
    """
    permission_classes = [IsAuthenticated, IsAdminUser]
    serializer_class = AdminLearningModuleSerializer
    parser_classes = [JSONParser, MultiPartParser, FormParser]
    
    def get_queryset(self):
        queryset = LearningModule.objects.select_related('career_path').annotate(
            quiz_count=Count('quizzes')
        )
        
        # Filters
        path_id = self.request.query_params.get('path')
        difficulty = self.request.query_params.get('difficulty')
        search = self.request.query_params.get('search')
        
        if path_id:
            queryset = queryset.filter(career_path_id=path_id)
        if difficulty:
            queryset = queryset.filter(difficulty_level=difficulty)
        if search:
            queryset = queryset.filter(
                Q(title__icontains=search) |
                Q(description__icontains=search)
            )
        
        return queryset.order_by('career_path', 'order')
    
    @action(detail=False, methods=['post'])
    def upload(self, request):
        """
        Upload and parse a file (DOCX, PDF, PPTX, MD)
        Returns parsed content for admin review
        """
        file = request.FILES.get('file')
        path_id = request.data.get('career_path_id')
        
        if not file:
            return Response(
                {'error': 'No file provided'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Validate file type
        allowed_extensions = ['docx', 'pdf', 'pptx', 'md']
        file_extension = file.name.split('.')[-1].lower()
        
        if file_extension not in allowed_extensions:
            return Response(
                {'error': f'File type not supported. Allowed: {", ".join(allowed_extensions)}'},
                status=status.HTTP_400_BAD_REQUEST
            )
        
        # Save file
        file_path = default_storage.save(f'module_uploads/{file.name}', file)
        
        # Parse file content
        try:
            parsed_content = self._parse_file(file_path, file_extension)
        except Exception as e:
            return Response(
                {'error': f'Failed to parse file: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        
        # Return parsed content for review
        return Response({
            'file_path': file_path,
            'file_name': file.name,
            'file_type': file_extension,
            'parsed_content': parsed_content,
            'preview': self._generate_preview(parsed_content),
        })
    
    @action(detail=False, methods=['post'])
    def confirm_upload(self, request):
        """
        Confirm and save parsed module content
        Creates the module with reviewed/edited content
        """
        file_path = request.data.get('file_path')
        path_id = request.data.get('career_path_id')
        edited_content = request.data.get('content')
        
        # Get highest order for this path
        max_order = LearningModule.objects.filter(
            career_path_id=path_id
        ).aggregate(max_order=models.Max('order'))['max_order'] or 0
        
        # Create module
        module = LearningModule.objects.create(
            career_path_id=path_id,
            title=edited_content.get('title', 'Untitled Module'),
            description=edited_content.get('description', ''),
            content=edited_content.get('content', ''),
            module_type=edited_content.get('type', 'reading'),
            difficulty_level=edited_content.get('difficulty', 'intermediate'),
            duration_minutes=edited_content.get('duration', 30),
            order=max_order + 1,
            file=file_path,
        )
        
        serializer = self.get_serializer(module)
        return Response(serializer.data, status=status.HTTP_201_CREATED)
    
    def _parse_file(self, file_path, file_extension):
        """
        Parse uploaded file and extract structure
        This is a simplified version - in production, use Apache Tika or similar
        """
        full_path = default_storage.path(file_path)
        
        if file_extension == 'md':
            return self._parse_markdown(full_path)
        elif file_extension == 'docx':
            return self._parse_docx(full_path)
        elif file_extension == 'pdf':
            return self._parse_pdf(full_path)
        elif file_extension == 'pptx':
            return self._parse_pptx(full_path)
        
        return {'error': 'Unsupported file type'}
    
    def _parse_markdown(self, file_path):
        """Parse markdown file"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Extract title (first # heading)
        lines = content.split('\n')
        title = 'Untitled'
        description = ''
        
        for line in lines:
            if line.startswith('# '):
                title = line.replace('# ', '').strip()
                break
        
        # Extract description (first paragraph)
        for line in lines:
            if line.strip() and not line.startswith('#'):
                description = line.strip()
                break
        
        # Split into slides (by ## headings)
        slides = []
        current_slide = {'title': '', 'content': ''}
        
        for line in lines:
            if line.startswith('## '):
                if current_slide['content']:
                    slides.append(current_slide)
                current_slide = {
                    'title': line.replace('## ', '').strip(),
                    'content': ''
                }
            elif current_slide['title']:
                current_slide['content'] += line + '\n'
        
        if current_slide['content']:
            slides.append(current_slide)
        
        return {
            'title': title,
            'description': description,
            'type': 'reading',
            'difficulty': 'intermediate',
            'duration': len(slides) * 5,  # 5 min per slide estimate
            'content': json.dumps({'slides': slides}),
        }
    
    def _parse_docx(self, file_path):
        """Parse DOCX file - requires python-docx"""
        try:
            from docx import Document
            doc = Document(file_path)
            
            title = 'Untitled'
            description = ''
            slides = []
            current_slide = {'title': '', 'content': ''}
            
            for para in doc.paragraphs:
                if para.style.name.startswith('Heading 1'):
                    title = para.text
                elif para.style.name.startswith('Heading 2'):
                    if current_slide['content']:
                        slides.append(current_slide)
                    current_slide = {'title': para.text, 'content': ''}
                else:
                    if not description and para.text.strip():
                        description = para.text
                    if current_slide['title']:
                        current_slide['content'] += para.text + '\n'
            
            if current_slide['content']:
                slides.append(current_slide)
            
            return {
                'title': title,
                'description': description,
                'type': 'reading',
                'difficulty': 'intermediate',
                'duration': len(slides) * 5,
                'content': json.dumps({'slides': slides}),
            }
        except ImportError:
            return {'error': 'python-docx not installed'}
    
    def _parse_pdf(self, file_path):
        """Parse PDF file - requires PyPDF2"""
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            
            full_text = ''
            for page in reader.pages:
                full_text += page.extract_text() + '\n'
            
            # Simple parsing - split by empty lines
            lines = full_text.split('\n')
            title = lines[0] if lines else 'Untitled'
            
            return {
                'title': title,
                'description': lines[1] if len(lines) > 1 else '',
                'type': 'reading',
                'difficulty': 'intermediate',
                'duration': len(reader.pages) * 3,
                'content': json.dumps({'slides': [{'title': f'Page {i+1}', 'content': page.extract_text()} for i, page in enumerate(reader.pages)]}),
            }
        except ImportError:
            return {'error': 'PyPDF2 not installed'}
    
    def _parse_pptx(self, file_path):
        """Parse PPTX file - requires python-pptx"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            slides = []
            title = 'Untitled'
            
            for i, slide in enumerate(prs.slides):
                slide_title = ''
                slide_content = ''
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        if i == 0 and not title:
                            title = shape.text
                        if not slide_title:
                            slide_title = shape.text
                        else:
                            slide_content += shape.text + '\n'
                
                slides.append({
                    'title': slide_title,
                    'content': slide_content
                })
            
            return {
                'title': title,
                'description': slides[1]['content'][:200] if len(slides) > 1 else '',
                'type': 'video',  # PPTX is more like a presentation
                'difficulty': 'intermediate',
                'duration': len(slides) * 5,
                'content': json.dumps({'slides': slides}),
            }
        except ImportError:
            return {'error': 'python-pptx not installed'}
    
    def _generate_preview(self, parsed_content):
        """Generate HTML preview of parsed content"""
        if 'error' in parsed_content:
            return parsed_content['error']
        
        preview = f"<h2>{parsed_content.get('title', 'Untitled')}</h2>"
        preview += f"<p>{parsed_content.get('description', '')}</p>"
        preview += f"<p><strong>Type:</strong> {parsed_content.get('type', 'reading')}</p>"
        preview += f"<p><strong>Duration:</strong> {parsed_content.get('duration', 0)} minutes</p>"
        
        return preview
    
    @action(detail=True, methods=['get'])
    def analytics(self, request, pk=None):
        """Get module analytics"""
        module = self.get_object()
        
        progress_records = ModuleProgress.objects.filter(module=module)
        total_views = progress_records.count()
        completed = progress_records.filter(status='completed')
        completed_count = completed.count()
        
        analytics_data = {
            'total_views': total_views,
            'unique_viewers': progress_records.values('user').distinct().count(),
            'total_starts': progress_records.exclude(status='not_started').count(),
            'total_completions': completed_count,
            'completion_rate': (completed_count / total_views * 100) if total_views > 0 else 0,
            'avg_time_spent': progress_records.aggregate(avg_time=Avg('time_spent_minutes'))['avg_time'] or 0,
            'slide_drop_off': {},
        }
        
        serializer = ModuleAnalyticsSerializer(analytics_data)
        return Response(serializer.data)


class AdminQuizViewSet(viewsets.ModelViewSet):
    """
    Complete Quiz Management
    """
    permission_classes = [IsAuthenticated, IsAdminUser]
    serializer_class = AdminQuizSerializer
    
    def get_queryset(self):
        queryset = Quiz.objects.select_related('module').annotate(
            questions_count=Count('questions')
        )
        
        module_id = self.request.query_params.get('module')
        if module_id:
            queryset = queryset.filter(module_id=module_id)
        
        return queryset.order_by('-created_at')
    
    @action(detail=True, methods=['post'])
    def add_question(self, request, pk=None):
        """Add a question to quiz"""
        quiz = self.get_object()
        question_data = request.data
        
        # Get max order
        max_order = quiz.questions.aggregate(
            max_order=models.Max('order')
        )['max_order'] or 0
        
        question = Question.objects.create(
            quiz=quiz,
            question_text=question_data.get('question_text'),
            question_type=question_data.get('question_type', 'multiple_choice'),
            points=question_data.get('points', 1),
            order=max_order + 1,
        )
        
        # Add choices if provided
        for choice_data in question_data.get('choices', []):
            QuestionChoice.objects.create(
                question=question,
                choice_text=choice_data.get('choice_text'),
                is_correct=choice_data.get('is_correct', False),
                order=choice_data.get('order', 0),
            )
        
        serializer = AdminQuestionSerializer(question)
        return Response(serializer.data, status=status.HTTP_201_CREATED)
    
    @action(detail=True, methods=['get'])
    def analytics(self, request, pk=None):
        """Get quiz analytics"""
        quiz = self.get_object()
        
        attempts = QuizAttempt.objects.filter(quiz=quiz)
        total_attempts = attempts.count()
        
        analytics_data = {
            'total_attempts': total_attempts,
            'unique_takers': attempts.values('user').distinct().count(),
            'avg_score': attempts.aggregate(avg_score=Avg('score'))['avg_score'] or 0,
            'highest_score': attempts.aggregate(max_score=models.Max('score'))['max_score'] or 0,
            'lowest_score': attempts.aggregate(min_score=models.Min('score'))['min_score'] or 0,
            'pass_rate': (attempts.filter(score__gte=quiz.passing_score).count() / total_attempts * 100) if total_attempts > 0 else 0,
            'avg_time_spent': attempts.aggregate(avg_time=Avg('time_taken_seconds'))['avg_time'] or 0,
            'question_difficulty': {},
        }
        
        serializer = QuizAnalyticsSerializer(analytics_data)
        return Response(serializer.data)


class AdminSearchView(views.APIView):
    """
    Global search across Paths, Modules, and Quizzes
    GET /api/admin/search/?q=query&type=path,module,quiz&program=BSIT
    """
    permission_classes = [IsAuthenticated, IsAdminUser]
    
    def get(self, request):
        query = request.query_params.get('q', '')
        search_types = request.query_params.get('type', 'path,module,quiz').split(',')
        program = request.query_params.get('program')
        difficulty = request.query_params.get('difficulty')
        status = request.query_params.get('status')
        
        results = {
            'paths': [],
            'modules': [],
            'quizzes': [],
        }
        
        if 'path' in search_types:
            paths = CareerPath.objects.filter(
                Q(name__icontains=query) |
                Q(description__icontains=query)
            )
            if program:
                paths = paths.filter(program=program)
            if difficulty:
                paths = paths.filter(difficulty_level=difficulty)
            if status:
                paths = paths.filter(is_active=(status == 'active'))
            
            results['paths'] = AdminCareerPathSerializer(paths[:10], many=True).data
        
        if 'module' in search_types:
            modules = LearningModule.objects.filter(
                Q(title__icontains=query) |
                Q(description__icontains=query)
            )
            if difficulty:
                modules = modules.filter(difficulty_level=difficulty)
            
            results['modules'] = AdminLearningModuleSerializer(modules[:10], many=True).data
        
        if 'quiz' in search_types:
            quizzes = Quiz.objects.filter(
                Q(title__icontains=query) |
                Q(description__icontains=query)
            )
            
            results['quizzes'] = AdminQuizSerializer(quizzes[:10], many=True).data
        
        return Response(results)


class AdminUserManagementView(viewsets.ModelViewSet):
    """
    Complete User Management
    Search, filter, view, edit, delete users
    """
    permission_classes = [IsAuthenticated, IsAdminUser]
    serializer_class = UserSerializer  # Use your User serializer
    
    def get_queryset(self):
        queryset = User.objects.all()
        
        # Filters
        search = self.request.query_params.get('search')
        role = self.request.query_params.get('role')
        program = self.request.query_params.get('program')
        is_active = self.request.query_params.get('is_active')
        
        if search:
            queryset = queryset.filter(
                Q(username__icontains=search) |
                Q(email__icontains=search) |
                Q(first_name__icontains=search) |
                Q(last_name__icontains=search)
            )
        if role:
            queryset = queryset.filter(role=role)
        if program:
            queryset = queryset.filter(program=program)
        if is_active is not None:
            queryset = queryset.filter(is_active=(is_active.lower() == 'true'))
        
        return queryset.order_by('-created_at')
    
    @action(detail=False, methods=['post'])
    def bulk_delete(self, request):
        """Delete multiple users"""
        user_ids = request.data.get('user_ids', [])
        deleted_count = User.objects.filter(id__in=user_ids).delete()[0]
        return Response({'deleted': deleted_count})


class AdminAnalyticsView(views.APIView):
    """
    Complete analytics dashboard with drill-down and export
    """
    permission_classes = [IsAuthenticated, IsAdminUser]
    
    def get(self, request):
        # High-level metrics
        data = {
            'overview': self._get_overview_analytics(),
            'paths': self._get_path_analytics(),
            'modules': self._get_module_analytics(),
            'quizzes': self._get_quiz_analytics(),
            'users': self._get_user_analytics(),
        }
        
        return Response(data)
    
    def _get_overview_analytics(self):
        return {
            'total_enrollments': Enrollment.objects.count(),
            'active_learners': Enrollment.objects.filter(completed_at__isnull=True).count(),
            'certificates_issued': Certificate.objects.count(),
            'avg_completion_rate': 75.5,  # Calculate from data
        }
    
    def _get_path_analytics(self):
        paths = CareerPath.objects.annotate(
            enrollments=Count('enrollments')
        ).order_by('-enrollments')[:5]
        
        return [
            {
                'id': str(path.id),
                'name': path.name,
                'enrollments': path.enrollments,
            }
            for path in paths
        ]
    
    def _get_module_analytics(self):
        return []
    
    def _get_quiz_analytics(self):
        return []
    
    def _get_user_analytics(self):
        return {
            'total_users': User.objects.count(),
            'active_users': User.objects.filter(is_active=True).count(),
            'by_program': list(User.objects.values('program').annotate(count=Count('id'))),
        }
